"""Sinh file slide PPTX bao cao du an Food Moo Duu."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
OUT_FILE = ROOT / "docs" / "food_moo_duu_slides.pptx"
IMG_DIR = ROOT / "docs" / "images"

sys.path.insert(0, str(ROOT / "scripts"))
from eval_report_utils import load_eval_bundle  # noqa: E402

# Bang mau
NAVY = RGBColor(0x1F, 0x2D, 0x3D)
ORANGE = RGBColor(0xE8, 0x6A, 0x17)
TEAL = RGBColor(0x12, 0x7C, 0x8A)
LIGHT = RGBColor(0xF5, 0xF1, 0xE8)
GREY = RGBColor(0x55, 0x5B, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Kich thuoc slide 16:9
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def add_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def style_run(run, size, color, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def add_accent_bar(slide, color=ORANGE):
    bar = slide.shapes.add_shape(
        1, Emu(0), Emu(0), Emu(160000), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def title_slide(prs, title, subtitle, members):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    add_accent_bar(slide, ORANGE)

    _, tf = add_box(slide, 900000, 1700000, 10400000, 1800000)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, 44, WHITE, bold=True)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = subtitle
    style_run(r2, 20, ORANGE, italic=True)

    _, tf3 = add_box(slide, 900000, 4100000, 10400000, 1800000)
    head = tf3.paragraphs[0]
    rh = head.add_run()
    rh.text = "Nhóm thực hiện"
    style_run(rh, 16, TEAL, bold=True)
    for name, role in members:
        pm = tf3.add_paragraph()
        rm = pm.add_run()
        rm.text = f"{name}  -  {role}"
        style_run(rm, 16, LIGHT)
    return slide


def content_slide(prs, title, blocks, owner=None):
    """blocks: list cua (loai, noi_dung).
    loai: 'h' (heading), 'b' (bullet cap 1), 's' (bullet cap 2), 'note'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_accent_bar(slide, ORANGE)

    # Tieu de
    _, tf_title = add_box(slide, 500000, 350000, 9800000, 900000)
    p = tf_title.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, 30, NAVY, bold=True)

    # Owner badge
    if owner:
        badge = slide.shapes.add_shape(1, Emu(9400000), Emu(380000), Emu(2400000), Emu(520000))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        btf = badge.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = owner
        style_run(br, 12, WHITE, bold=True)

    # Noi dung
    _, tf = add_box(slide, 600000, 1450000, 11000000, 5000000)
    first = True
    for kind, text in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if kind == "h":
            p.space_before = Pt(10)
            r = p.add_run()
            r.text = text
            style_run(r, 19, ORANGE, bold=True)
        elif kind == "b":
            p.space_before = Pt(5)
            r = p.add_run()
            r.text = "•  " + text
            style_run(r, 16, NAVY)
        elif kind == "s":
            p.level = 1
            r = p.add_run()
            r.text = "–  " + text
            style_run(r, 14, GREY)
        elif kind == "note":
            p.space_before = Pt(8)
            r = p.add_run()
            r.text = text
            style_run(r, 13, TEAL, italic=True)
    return slide


def section_slide(prs, title, owner):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, TEAL)
    add_accent_bar(slide, ORANGE)
    _, tf = add_box(slide, 900000, 2600000, 10400000, 1600000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, 36, WHITE, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"Phụ trách: {owner}"
    style_run(r2, 18, LIGHT, italic=True)
    return slide


def add_image_centered(slide, image_name, top, max_w, max_h):
    """Chen anh, can giua ngang, giu ti le, gioi han trong khung max_w x max_h (EMU)."""
    path = IMG_DIR / image_name
    if not path.exists():
        return
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    width = max_w
    height = int(width / ratio)
    if height > max_h:
        height = max_h
        width = int(height * ratio)
    left = int((SLIDE_W - width) / 2)
    slide.shapes.add_picture(str(path), Emu(left), Emu(top), Emu(width), Emu(height))


def image_slide(prs, title, image_name, blocks, owner=None, img_height=3400000):
    """Slide: tieu de + bullet ngan o tren + hinh minh hoa lon o duoi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_accent_bar(slide, ORANGE)

    _, tf_title = add_box(slide, 500000, 300000, 9800000, 850000)
    p = tf_title.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, 26, NAVY, bold=True)

    if owner:
        badge = slide.shapes.add_shape(1, Emu(9400000), Emu(330000), Emu(2400000), Emu(520000))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = owner
        style_run(br, 12, WHITE, bold=True)

    # Bullet ngan
    _, tf = add_box(slide, 600000, 1250000, 11000000, 1900000)
    first = True
    for kind, text in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if kind == "h":
            r = p.add_run()
            r.text = text
            style_run(r, 15, ORANGE, bold=True)
        elif kind == "b":
            r = p.add_run()
            r.text = "•  " + text
            style_run(r, 13, NAVY)
        elif kind == "note":
            r = p.add_run()
            r.text = text
            style_run(r, 12, TEAL, italic=True)

    # Hinh minh hoa o duoi
    add_image_centered(slide, image_name, top=3300000, max_w=10600000, max_h=img_height)
    return slide


def build():
    run_id, eval_bundle, _ = load_eval_bundle()
    summary = eval_bundle["summary"]
    l1_no_rl = eval_bundle["layer1_no_rl"]["summary"]
    l1_with_rl = eval_bundle["layer1_with_rl"]["summary"]
    l2_oracle = eval_bundle["layer2_oracle"]["metrics"]
    l2_behavioral = eval_bundle["layer2_behavioral"]["metrics"]
    l3 = eval_bundle["layer3"]
    behavioral_n = int(eval_bundle["layer2_behavioral"].get("samples", 0) or 0)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    members = [
        ("Trần Việt Hoàng", "Layer 1 - Định nghĩa Tag + Kiến trúc tổng thể"),
        ("Bùi Đức Anh", "Layer 2 - Adaptive Recommendation"),
        ("Trần Bình Minh", "Layer 3 - Genetic Response"),
    ]

    # Slide 1 - Bia
    title_slide(
        prs,
        "Food Moo Duu - Hệ thống gợi ý món ăn thông minh offline",
        '"Là bò, sáng phải rống - Là người, sống phải ráng"',
        members,
    )

    # Slide 2 - Van de & muc tieu
    content_slide(
        prs,
        "Vấn đề & Mục tiêu",
        [
            ("h", "Vấn đề"),
            ("b", "Người dùng mô tả nhu cầu ăn uống bằng ngôn ngữ tự nhiên, đa chiều (thời gian, thời tiết, tâm trạng, khẩu vị)"),
            ("b", "Cần hệ thống offline, không phụ thuộc API cloud"),
            ("b", "Cần học từ phản hồi khi người dùng chọn hoặc bỏ qua gợi ý"),
            ("h", "Mục tiêu"),
            ("b", "Xây dựng chatbot gợi ý món ăn theo kiến trúc 3 tầng"),
            ("b", "Theo dõi ngữ cảnh hội thoại nhiều lượt (multi-turn)"),
            ("b", "Tự cập nhật mô hình gợi ý và câu trả lời theo feedback ngầm"),
        ],
        owner="Hoàng",
    )

    # Slide 3 - Tong quan giai phap
    content_slide(
        prs,
        "Tổng quan giải pháp",
        [
            ("h", "Ý tưởng cốt lõi"),
            ("b", "Pipeline xử lý tuần tự 3 layer: DL + học trực tuyến + thuật toán di truyền"),
            ("h", "Điểm nổi bật"),
            ("b", "Offline end-to-end qua main.py -> FoodSuggestionPipeline"),
            ("b", "Mỗi layer chạy độc lập được (train/test/docker riêng)"),
            ("b", "Dữ liệu train và runtime tách biệt rõ ràng"),
            ("h", "Công nghệ"),
            ("b", "Python, PyTorch, scikit-learn, pandas, Docker, Makefile"),
        ],
        owner="Hoàng",
    )

    # Slide 4 - Kien truc he thong (HINH BAT BUOC)
    image_slide(
        prs,
        "Kiến trúc hệ thống (sơ đồ tổng thể)",
        "system_architecture.png",
        [
            ("h", "Luồng xử lý + 2 vòng phản hồi"),
            ("b", "User -> L1 (Intent+DST) -> L2 (Scoring) -> L3 (Response) -> User"),
            ("b", "Feedback: Hebbian update (L2) + Fitness update (L3); RL offline cho L1"),
        ],
        owner="Hoàng",
    )

    # ===== PHAN MUC LAYER 1 =====
    section_slide(prs, "Layer 1 - Intent & Dialog State Tracking", "Trần Việt Hoàng")

    # Slide 6 - Layer 1: Y tuong & logic (HINH)
    image_slide(
        prs,
        "Layer 1: Ý tưởng & logic thực hiện",
        "layer1_logic.png",
        [
            ("h", "Ý tưởng: biến câu nói -> bộ tag ngữ cảnh (multi-label) + học embedding theo ngữ cảnh"),
            ("b", "Chuẩn hóa VN + tokenize -> Encoder vi-SBERT (mean pooling)"),
            ("b", "2 head: Classification (sigmoid -> ~53 tag) + Projection (metric learning)"),
            ("note", "Loss = BCE multi-label + alpha(0.35) * metric loss (overlap-based similarity)"),
        ],
        owner="Hoàng",
    )

    # Slide 7 - Layer 1: DST logic
    content_slide(
        prs,
        "Layer 1: Dialog State Tracking (logic DST)",
        [
            ("h", "Vai trò: tích lũy ngữ cảnh nhiều lượt, xử lý người dùng đổi ý"),
            ("h", "3 quy luật cập nhật state (update_context)"),
            ("b", "1. Decay: score *= 0.55 mỗi lượt (session state yếu dần, ưu tiên lượt mới)"),
            ("b", "2. Accumulation: score += 0.88 * confidence (raw tag lượt hiện tại ảnh hưởng mạnh)"),
            ("b", "3. Conflict resolution: cặp tag đối nghịch -> giảm bên yếu theo gap (beta=0.4)"),
            ("note", "Cấu hình: data/common_config.json -> dst_hyperparameters"),
            ("s", "VD: weather_hot <-> weather_cold, pref_spicy <-> pref_bland"),
            ("note", "Output: context_scores (clamp [0,1]) -> Layer 2 & Layer 3; lưu session_state.json"),
        ],
        owner="Hoàng",
    )

    # Slide 8 - Du lieu & huan luyen Layer 1
    content_slide(
        prs,
        "Dữ liệu & Huấn luyện Layer 1",
        [
            ("h", "Dataset"),
            ("b", "5 phiên bản: dataset_v001 -> dataset_v005, mỗi epoch ~160 mẫu"),
            ("b", "Train trên active_dataset hoặc --all-datasets; gộp thêm RL offline"),
            ("h", "Cấu hình (dl_config.json)"),
            ("b", "8 epochs, batch 32, lr 0.001, threshold tag 0.3"),
            ("b", "Artifacts: intent_model.pt, vocab.json, intent_model_meta.json"),
            ("b", "Đánh giá: micro/macro F1 trên validation (20%)"),
            ("note", "Lệnh: make layer1-train-active | make layer1-chat"),
        ],
        owner="Hoàng",
    )

    # ===== PHAN MUC LAYER 2 =====
    section_slide(prs, "Layer 2 - Adaptive Recommendation", "Bùi Đức Anh")

    # Slide 10 - Layer 2: Y tuong & logic (HINH)
    image_slide(
        prs,
        "Layer 2: Ý tưởng & logic thực hiện",
        "layer2_logic.png",
        [
            ("h", "Ý tưởng: mỗi món = vector trọng số trên ~53 tag; gợi ý = so khớp ngữ cảnh"),
            ("b", "Mở rộng context qua similarity tag (layer2_config.json)"),
            ("b", "Linear scoring: score(món) = SUM activation(tag) * weight(món, tag)"),
            ("note", "Xếp hạng theo (score, popularity, id) -> Top-K; nguồn: 120 món (9 đồ uống)"),
        ],
        owner="Đức Anh",
    )

    # Slide 11 - Layer 2: Hebbian logic
    content_slide(
        prs,
        "Layer 2: Hebbian Online Learning (logic cập nhật)",
        [
            ("h", "Reward - món được chọn (apply_context_to_dish)"),
            ("b", "Tag active (>= 0.25): weight += lr_positive(0.08) * activation"),
            ("b", "Tag không active: weight -= penalty * (0.25 - activation)"),
            ("h", "Penalty - món trong Top-K nhưng không được chọn"),
            ("b", "weight -= lr_negative(0.02) * activation"),
            ("b", "Trọng số luôn clamp [-1, 1] rồi save() xuống runtime matrix"),
            ("note", "Demo: hiển thị score_before, score_after, delta để thấy hệ thống đã học"),
        ],
        owner="Đức Anh",
    )

    # Slide 12 - Quy uoc du lieu & van hanh Layer 2
    content_slide(
        prs,
        "Quy ước dữ liệu & Vận hành Layer 2",
        [
            ("h", "Convention"),
            ("b", "Mỗi layer có datasets.json khai báo active_dataset"),
            ("b", "Data train tách khỏi runtime (dishes_runtime.json, feedback_reports.jsonl)"),
            ("h", "Tiện ích vận hành"),
            ("b", "make layer2-migrate (canonical) | layer2-check (schema drift)"),
            ("b", "make layer2-reset-runtime | layer2-test"),
        ],
        owner="Đức Anh",
    )

    # ===== PHAN MUC LAYER 3 =====
    section_slide(prs, "Layer 3 - Genetic Response", "Trần Bình Minh")

    # Slide 14 - Layer 3: Y tuong & logic (HINH)
    image_slide(
        prs,
        "Layer 3: Ý tưởng & logic thực hiện",
        "layer3_logic.png",
        [
            ("h", "Ý tưởng: câu thoại được 'tiến hóa' từ gene pool theo mood; câu tốt -> fitness cao"),
            ("b", "Suy mood_key từ context_scores -> khởi tạo 8 chromosome (bộ 3 gene)"),
            ("b", "Chọn cá thể: epsilon-greedy (0.2) hoặc roulette theo fitness; mutation slang"),
            ("note", "Ghép opening+action+closing -> câu trả lời, chèn tên món vào {foods}"),
        ],
        owner="Minh",
    )

    # Slide 15 - Layer 3: Fitness update logic
    content_slide(
        prs,
        "Layer 3: Fitness Update (logic phản hồi ngầm)",
        [
            ("h", "Cập nhật fitness theo implicit feedback (update_fitness)"),
            ("b", "Người dùng chọn món -> fitness += 0.25 (củng cố chromosome tốt)"),
            ("b", "Bỏ qua / thoát -> fitness -= 0.2 (tối thiểu 0.05 để không triệt tiêu)"),
            ("h", "Lưu trữ"),
            ("b", "Ghi history + runtime_fitness riêng vào fitness_history.json"),
            ("note", "Fitness cao -> xác suất roulette chọn cao hơn -> câu thoại tự tối ưu theo thời gian"),
        ],
        owner="Minh",
    )

    # Slide 16 - Luong runtime end-to-end
    content_slide(
        prs,
        "Luồng runtime end-to-end",
        [
            ("h", "Kịch bản demo (main.py)"),
            ("b", "Chat lượt 1: 'Tối nay trời lạnh, tôi muốn món nước nóng'"),
            ("b", "Chat lượt 2..5: có thể quay xe đổi ý (DST cập nhật ngữ cảnh)"),
            ("b", "Sau mỗi lượt: hệ thống trả câu trả lời (L3) + top-K gợi ý (L2)"),
            ("b", "User chọn món -> Hebbian (L2) + fitness (L3); không chọn -> chỉ fitness"),
            ("note", "Entry: process_turn() -> apply_feedback() / apply_abandon_feedback()"),
        ],
        owner="Hoàng",
    )

    # Slide 17 - RL Layer 1 offline
    content_slide(
        prs,
        "Reinforcement Learning cho Layer 1 (offline)",
        [
            ("h", "Mục tiêu: thu thập feedback chọn món để train tăng cường - không train realtime"),
            ("b", "1. Logging: ghi event khi user chọn món -> selected_dish_events.jsonl"),
            ("b", "2. Simulator: sinh dữ liệu giả lập -> selected_dish_events_simulated.jsonl"),
            ("b", "3. Offline train: chuyển feedback -> intent_train_data_rl.json"),
            ("b", "Thống kê: 102 events -> 102 train samples (stats.json)"),
            ("note", "Lệnh: make layer1-rl-generate -> layer1-rl-train -> layer1-rl-check"),
        ],
        owner="Hoàng",
    )

    # ===== PHAN MUC DANH GIA =====
    section_slide(prs, "Đánh giá hiệu quả", "Cả nhóm")

    # Slide 19 - Pipeline danh gia
    content_slide(
        prs,
        "Pipeline đánh giá & phương pháp",
        [
            ("h", "Package: src/evaluation/ (metrics, layer1/2/3_eval, pipeline_eval)"),
            ("b", "make layer3-simulate — bổ sung data giả lập L3"),
            ("b", "make eval-run — đánh giá toàn bộ (skip train L1)"),
            ("b", "make eval-train-and-run — train ablation L1 + đánh giá"),
            ("h", "Tập dữ liệu đánh giá"),
            ("b", "L1: validation 20% intent (all datasets) — 142 mẫu"),
            ("b", "L2 Oracle: intent_samples.csv tag lý tưởng — 709 mẫu"),
            ("b", f"L2 Behavioral: RL events — {behavioral_n} events"),
            ("b", "L3: fitness_history (12 runtime + 50 simulated) — 62 lượt"),
            ("b", "DST runtime: decay=0.55, alpha=0.88, beta=0.4 (ưu tiên raw tag)"),
            ("note", "Metric: F1 (L1), Hit@K/MRR/NDCG (L2), success rate/fitness (L3), feedback delta"),
        ],
        owner="Hoàng",
    )

    # Slide 20 - Ket qua tong hop (HINH)
    image_slide(
        prs,
        "Kết quả đánh giá tổng hợp",
        "evaluation_metrics.png",
        [
            ("h", f"Run: {run_id} | DST: decay=0.55, alpha=0.88, beta=0.4"),
            ("b", f"L1 Macro F1 = {summary['layer1_with_rl_macro_f1']} | L2 Oracle Hit@5 = {summary['layer2_oracle_hit_at_5']} | L2 Behavioral Hit@5 = {summary['layer2_behavioral_hit_at_5']}"),
            ("b", f"L3 Success rate = {round(summary['layer3_success_rate']*100, 1)}% ({int(l3.get('total_updates', 0))} lượt) | E2E Hit@5 = {summary['pipeline_e2e_with_rl_hit_at_5']} | Feedback delta = {summary['pipeline_feedback_delta_mean']:+.3f}"),
            ("note", "Oracle cao -> scoring tốt khi tag chuẩn; behavioral thấp -> context runtime khó hơn"),
        ],
        owner="Cả nhóm",
        img_height=3600000,
    )

    # Slide 21 - Chi tiet metric tung layer
    content_slide(
        prs,
        "Chi tiết metric từng layer",
        [
            ("h", f"Layer 1 (val {int(l1_no_rl['samples'])})"),
            ("b", f"Micro F1 = {l1_no_rl['micro_f1']}, Macro F1 = {l1_no_rl['macro_f1']}, Precision = {l1_no_rl['micro_precision']}, Recall = {l1_no_rl['micro_recall']}"),
            ("b", f"Ablation RL: delta macro F1 = {summary['layer1_delta_macro_f1']}"),
            ("h", "Layer 2"),
            ("b", f"Oracle (709): Hit@5 = {l2_oracle.get('hit_at_5')}, NDCG@5 = {round(l2_oracle.get('ndcg_at_5', 0), 3)} — upper bound scoring"),
            ("b", f"Behavioral ({behavioral_n}): Hit@5 = {l2_behavioral.get('hit_at_5')}, MRR = {l2_behavioral.get('mrr')}, mean rank ~ {round(l2_behavioral.get('mean_rank', 0), 2)}"),
            ("h", "Layer 3 & Pipeline"),
            ("b", f"L3: success {round(summary['layer3_success_rate']*100, 1)}% (runtime {round(l3.get('runtime_success_rate', 0)*100, 1)}%, simulated {round(l3.get('simulated_success_rate', 0)*100, 1)}%), {int(l3.get('unique_chromosomes', 0))} chromosome, avg fitness {l3.get('avg_fitness')}"),
            ("b", f"E2E Hit@5 = {summary['pipeline_e2e_with_rl_hit_at_5']}; Hebbian feedback delta mean {summary['pipeline_feedback_delta_mean']:+.3f}"),
            ("note", "Menu 120 mon da gan lai tag; runtime matrix da reset sau cap nhat catalog"),
        ],
        owner="Cả nhóm",
    )

    # Slide 22 - DevOps
    content_slide(
        prs,
        "Triển khai & DevOps",
        [
            ("h", "Makefile - một lệnh cho mỗi tác vụ"),
            ("b", "make setup | app-run | app-demo | eval-run | layer{1,2,3}-*"),
            ("h", "Docker"),
            ("b", "docker/layer1|2|3/docker-compose.yml + docker/app/ (full app)"),
            ("h", "Testing"),
            ("b", "test_layer2.py, test_layer2_integration.py, test_feedback_logging.py, test_metrics.py"),
        ],
        owner="Đức Anh",
    )

    # Slide 23 - Demo
    content_slide(
        prs,
        "Demo minh họa",
        [
            ("h", "Ví dụ non-interactive (make app-demo)"),
            ("s", "Chat 1: 'Tối nay trời lạnh, tôi muốn món nước nóng'"),
            ("s", "Bot: câu trả lời genetic + Top 5 món + score"),
            ("s", "Chat 2..5: 'Quay xe, giờ tôi muốn món nhanh và tiện lợi'"),
            ("s", "Feedback: chosen=dish_001 -> delta score"),
            ("h", "Điểm cần nhấn"),
            ("b", "DST xử lý 'quay xe' | score món đổi sau feedback | câu trả lời theo mood"),
        ],
        owner="Minh",
    )

    # Slide 24 - Han che & ket luan
    content_slide(
        prs,
        "Hạn chế & Kết luận",
        [
            ("h", "Ưu điểm"),
            ("b", "Modular, offline, DL multi-label + DST + Hebbian + GA + RL offline"),
            ("b", "Đã có pipeline đánh giá định lượng (make eval-run)"),
            ("h", "Kết quả nổi bật"),
            ("b", "DST tinh chỉnh: raw tag mạnh hơn (decay 0.55, alpha 0.88)"),
            ("b", f"L2 Oracle Hit@5 = {summary['layer2_oracle_hit_at_5']} | L3 success rate ~ {round(summary['layer3_success_rate']*100, 1)}% | Menu 120 mon da da dang"),
            ("h", "Hạn chế"),
            ("b", f"Dataset nhỏ; L2 behavioral Hit@5 thấp ({summary['layer2_behavioral_hit_at_5']}); L1 recall thấp ({round(l1_no_rl['micro_recall'], 2)})"),
            ("b", "L2/L3 eval chủ yếu simulated; chưa có UI"),
            ("h", "Hướng phát triển"),
            ("b", "Thu thập phiên chat thật, mở rộng corpus, UI, ablation RL L1"),
            ("note", "Food Moo Duu: pipeline NLP end-to-end tự học từ lựa chọn người dùng. Cảm ơn - Q&A"),
        ],
        owner="Cả nhóm",
    )

    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(f"Đã tạo slide: {OUT_FILE}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
